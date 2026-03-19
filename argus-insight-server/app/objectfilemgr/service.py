"""Object file manager service.

Implements S3-compatible object storage operations for both AWS S3 and MinIO.

Capabilities:
  1. Object CRUD: PutObject, GetObject, DeleteObject, CopyObject, HeadObject, DeleteObjects
  2. Directory (Prefix) management: ListObjectsV2, CreateFolder
  3. Large file support: Multipart Upload (initiate, sign parts, complete, abort)
  4. S3 Select: SQL query on CSV/JSON/Parquet objects
  5. Object Tagging: Get/Put/Delete tags
  6. File Browser Configuration: Dynamic configuration from database
"""

import logging
from collections import defaultdict
from datetime import UTC, datetime

from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession

from app.core.s3 import get_s3_client, get_s3_settings
from app.objectfilemgr.models import (
    ArgusConfigurationFilebrowser,
    ArgusConfigurationFilebrowserExtension,
    ArgusConfigurationFilebrowserPreview,
)
from app.usermgr.models import ArgusUser
from app.objectfilemgr.schemas import (
    BucketInfo,
    BucketListResponse,
    CompletedPart,
    CompleteMultipartResponse,
    CopyObjectResponse,
    CreateFolderResponse,
    DeleteObjectsResponse,
    DocumentPreviewResponse,
    EnsureUserBucketsResponse,
    FilebrowserConfigResponse,
    FolderInfo,
    HeadObjectResponse,
    ListObjectsResponse,
    MultipartUploadInitResponse,
    MultipartUploadPartUrl,
    MultipartUploadUrlsResponse,
    ObjectInfo,
    PresignedUploadUrlResponse,
    PresignedUrlResponse,
    PreviewCategoryResponse,
    S3SelectResponse,
    TablePreviewResponse,
    Tag,
    TaggingResponse,
)

logger = logging.getLogger(__name__)


# =========================================================================== #
# Bucket management
# =========================================================================== #


async def list_buckets() -> BucketListResponse:
    """List all S3 buckets."""
    async with get_s3_client() as s3:
        resp = await s3.list_buckets()
    buckets = []
    for b in resp.get("Buckets", []):
        creation_date = b.get("CreationDate")
        buckets.append(BucketInfo(
            name=b["Name"],
            creation_date=creation_date.isoformat() if creation_date else None,
        ))
    return BucketListResponse(buckets=buckets)


async def ensure_user_buckets(db: AsyncSession) -> EnsureUserBucketsResponse:
    """Ensure a user-<username> bucket exists for every user."""
    result = await db.execute(select(ArgusUser.username))
    usernames = [row[0] for row in result.all()]

    async with get_s3_client() as s3:
        resp = await s3.list_buckets()
        existing_buckets = {b["Name"] for b in resp.get("Buckets", [])}

        created: list[str] = []
        existing: list[str] = []

        for username in usernames:
            bucket_name = f"user-{username}"
            if bucket_name in existing_buckets:
                existing.append(bucket_name)
            else:
                await s3.create_bucket(Bucket=bucket_name)
                created.append(bucket_name)
                logger.info("Created user bucket: %s", bucket_name)

    return EnsureUserBucketsResponse(created=created, existing=existing)


def _format_dt(dt: datetime | str | None) -> str:
    """Convert a datetime to ISO 8601 string."""
    if dt is None:
        return ""
    if isinstance(dt, str):
        return dt
    return dt.astimezone(UTC).isoformat()


def _extract_name(prefix: str) -> str:
    """Extract the last segment of a prefix path as its display name."""
    trimmed = prefix.rstrip("/")
    return trimmed.rsplit("/", 1)[-1] if "/" in trimmed else trimmed


# =========================================================================== #
# 1. Object management
# =========================================================================== #


async def head_object(bucket: str, key: str) -> HeadObjectResponse:
    """Retrieve object metadata without downloading the object body."""
    async with get_s3_client() as s3:
        resp = await s3.head_object(Bucket=bucket, Key=key)
    return HeadObjectResponse(
        key=key,
        size=resp.get("ContentLength", 0),
        last_modified=_format_dt(resp.get("LastModified")),
        etag=resp.get("ETag", "").strip('"'),
        content_type=resp.get("ContentType"),
        storage_class=resp.get("StorageClass"),
        metadata=resp.get("Metadata", {}),
    )


async def get_object(bucket: str, key: str) -> dict:
    """Download an object. Returns the streaming body and metadata.

    The caller (router) is responsible for streaming the body to the HTTP response.
    """
    async with get_s3_client() as s3:
        resp = await s3.get_object(Bucket=bucket, Key=key)
        body = await resp["Body"].read()
    return {
        "body": body,
        "content_type": resp.get("ContentType", "application/octet-stream"),
        "content_length": resp.get("ContentLength", 0),
        "etag": resp.get("ETag", "").strip('"'),
        "last_modified": _format_dt(resp.get("LastModified")),
    }


async def put_object(
    bucket: str,
    key: str,
    body: bytes,
    content_type: str = "application/octet-stream",
) -> dict:
    """Upload a single object."""
    async with get_s3_client() as s3:
        resp = await s3.put_object(
            Bucket=bucket,
            Key=key,
            Body=body,
            ContentType=content_type,
        )
    etag = resp.get("ETag", "").strip('"')
    logger.info("PutObject: bucket=%s key=%s size=%d", bucket, key, len(body))
    return {"key": key, "etag": etag}


async def put_object_stream(
    bucket: str,
    key: str,
    file,
    content_type: str = "application/octet-stream",
    multipart_threshold: int = 8 * 1024 * 1024,
    chunk_size: int = 8 * 1024 * 1024,
) -> dict:
    """Upload a file using streaming. Uses multipart upload for large files.

    For files smaller than *multipart_threshold* the content is read into memory
    and uploaded with a single PutObject call. For larger files the content is
    streamed in *chunk_size* chunks via multipart upload so that only one chunk
    needs to be in memory at any time.
    """
    # Read first chunk to decide strategy
    first_chunk = await file.read(multipart_threshold + 1)

    if len(first_chunk) <= multipart_threshold:
        # Small file: single PutObject
        async with get_s3_client() as s3:
            resp = await s3.put_object(
                Bucket=bucket, Key=key, Body=first_chunk, ContentType=content_type,
            )
        etag = resp.get("ETag", "").strip('"')
        logger.info("PutObject: bucket=%s key=%s size=%d", bucket, key, len(first_chunk))
        return {"key": key, "etag": etag}

    # Large file: multipart upload
    async with get_s3_client() as s3:
        mpu = await s3.create_multipart_upload(
            Bucket=bucket, Key=key, ContentType=content_type,
        )
        upload_id = mpu["UploadId"]
        parts: list[dict] = []
        part_number = 1
        total_size = 0

        try:
            # Upload first chunk (already read, may be > threshold)
            offset = 0
            while offset < len(first_chunk):
                end = min(offset + chunk_size, len(first_chunk))
                chunk = first_chunk[offset:end]
                resp = await s3.upload_part(
                    Bucket=bucket, Key=key, UploadId=upload_id,
                    PartNumber=part_number, Body=chunk,
                )
                parts.append({
                    "PartNumber": part_number,
                    "ETag": resp["ETag"],
                })
                total_size += len(chunk)
                part_number += 1
                offset = end

            # Stream remaining chunks
            while True:
                chunk = await file.read(chunk_size)
                if not chunk:
                    break
                resp = await s3.upload_part(
                    Bucket=bucket, Key=key, UploadId=upload_id,
                    PartNumber=part_number, Body=chunk,
                )
                parts.append({
                    "PartNumber": part_number,
                    "ETag": resp["ETag"],
                })
                total_size += len(chunk)
                part_number += 1

            result = await s3.complete_multipart_upload(
                Bucket=bucket, Key=key, UploadId=upload_id,
                MultipartUpload={"Parts": parts},
            )
            etag = result.get("ETag", "").strip('"')
            logger.info(
                "PutObject(multipart): bucket=%s key=%s size=%d parts=%d",
                bucket, key, total_size, len(parts),
            )
            return {"key": key, "etag": etag}

        except Exception:
            await s3.abort_multipart_upload(
                Bucket=bucket, Key=key, UploadId=upload_id,
            )
            raise


async def delete_object(bucket: str, key: str) -> None:
    """Delete a single object."""
    async with get_s3_client() as s3:
        await s3.delete_object(Bucket=bucket, Key=key)
    logger.info("DeleteObject: bucket=%s key=%s", bucket, key)


async def delete_objects(bucket: str, keys: list[str]) -> DeleteObjectsResponse:
    """Delete multiple objects in a single request (up to 1000)."""
    async with get_s3_client() as s3:
        resp = await s3.delete_objects(
            Bucket=bucket,
            Delete={"Objects": [{"Key": k} for k in keys], "Quiet": False},
        )
    deleted = [d["Key"] for d in resp.get("Deleted", [])]
    errors = [
        {"key": e["Key"], "code": e.get("Code"), "message": e.get("Message")}
        for e in resp.get("Errors", [])
    ]
    logger.info("DeleteObjects: bucket=%s requested=%d deleted=%d", bucket, len(keys), len(deleted))
    return DeleteObjectsResponse(deleted=deleted, errors=errors)


async def copy_object(
    bucket: str,
    source_key: str,
    destination_key: str,
    source_bucket: str | None = None,
) -> CopyObjectResponse:
    """Copy an object within or across buckets.

    S3 has no native Move operation; move = copy + delete.
    """
    src_bucket = source_bucket or bucket
    copy_source = f"{src_bucket}/{source_key}"
    async with get_s3_client() as s3:
        resp = await s3.copy_object(
            Bucket=bucket,
            CopySource=copy_source,
            Key=destination_key,
        )
    etag = resp.get("CopyObjectResult", {}).get("ETag", "").strip('"')
    logger.info("CopyObject: %s/%s -> %s/%s", src_bucket, source_key, bucket, destination_key)
    return CopyObjectResponse(key=destination_key, etag=etag)


# =========================================================================== #
# 2. Directory (Prefix) management
# =========================================================================== #


async def list_objects(
    bucket: str,
    prefix: str = "",
    delimiter: str = "/",
    continuation_token: str | None = None,
    max_keys: int = 1000,
) -> ListObjectsResponse:
    """List objects using ListObjectsV2 with Delimiter support.

    When delimiter is "/", CommonPrefixes returns virtual folder names,
    and Contents returns only direct-child objects under the prefix.
    """
    params: dict = {
        "Bucket": bucket,
        "Prefix": prefix,
        "Delimiter": delimiter,
        "MaxKeys": max_keys,
    }
    if continuation_token:
        params["ContinuationToken"] = continuation_token

    async with get_s3_client() as s3:
        resp = await s3.list_objects_v2(**params)

    folders = [
        FolderInfo(prefix=cp["Prefix"], name=_extract_name(cp["Prefix"]))
        for cp in resp.get("CommonPrefixes", [])
    ]

    objects = [
        ObjectInfo(
            key=obj["Key"],
            size=obj.get("Size", 0),
            last_modified=_format_dt(obj.get("LastModified")),
            etag=obj.get("ETag", "").strip('"'),
            storage_class=obj.get("StorageClass"),
        )
        for obj in resp.get("Contents", [])
        # Skip the prefix marker itself (0-byte folder placeholder)
        if obj["Key"] != prefix
    ]

    return ListObjectsResponse(
        folders=folders,
        objects=objects,
        next_continuation_token=resp.get("NextContinuationToken"),
        is_truncated=resp.get("IsTruncated", False),
        key_count=resp.get("KeyCount", len(folders) + len(objects)),
    )


async def create_folder(bucket: str, key: str) -> CreateFolderResponse:
    """Create a virtual folder by putting a 0-byte object with a trailing slash."""
    folder_key = key if key.endswith("/") else key + "/"
    async with get_s3_client() as s3:
        await s3.put_object(Bucket=bucket, Key=folder_key, Body=b"")
    logger.info("CreateFolder: bucket=%s key=%s", bucket, folder_key)
    return CreateFolderResponse(key=folder_key)


# =========================================================================== #
# 3. Multipart upload
# =========================================================================== #


async def initiate_multipart_upload(
    bucket: str,
    key: str,
    content_type: str = "application/octet-stream",
) -> MultipartUploadInitResponse:
    """Initiate a multipart upload and return the upload_id."""
    async with get_s3_client() as s3:
        resp = await s3.create_multipart_upload(
            Bucket=bucket,
            Key=key,
            ContentType=content_type,
        )
    upload_id = resp["UploadId"]
    logger.info("InitiateMultipartUpload: bucket=%s key=%s upload_id=%s", bucket, key, upload_id)
    return MultipartUploadInitResponse(upload_id=upload_id, key=key)


async def _presigned_url_expiry() -> int:
    """Get the presigned URL expiry from DB settings."""
    cfg = await get_s3_settings()
    return int(cfg.get("object_storage_presigned_url_expiry", "3600"))


async def get_multipart_upload_urls(
    bucket: str,
    key: str,
    upload_id: str,
    part_numbers: list[int],
) -> MultipartUploadUrlsResponse:
    """Generate presigned URLs for uploading individual parts."""
    expiry = await _presigned_url_expiry()
    parts: list[MultipartUploadPartUrl] = []
    async with get_s3_client() as s3:
        for pn in part_numbers:
            url = await s3.generate_presigned_url(
                "upload_part",
                Params={
                    "Bucket": bucket,
                    "Key": key,
                    "UploadId": upload_id,
                    "PartNumber": pn,
                },
                ExpiresIn=expiry,
            )
            parts.append(MultipartUploadPartUrl(part_number=pn, url=url))
    return MultipartUploadUrlsResponse(parts=parts)


async def complete_multipart_upload(
    bucket: str,
    key: str,
    upload_id: str,
    parts: list[CompletedPart],
) -> CompleteMultipartResponse:
    """Complete a multipart upload by assembling the parts."""
    multipart = {
        "Parts": [
            {"PartNumber": p.part_number, "ETag": p.etag}
            for p in sorted(parts, key=lambda x: x.part_number)
        ]
    }
    async with get_s3_client() as s3:
        resp = await s3.complete_multipart_upload(
            Bucket=bucket,
            Key=key,
            UploadId=upload_id,
            MultipartUpload=multipart,
        )
    etag = resp.get("ETag", "").strip('"')
    logger.info(
        "CompleteMultipartUpload: bucket=%s key=%s upload_id=%s parts=%d",
        bucket, key, upload_id, len(parts),
    )
    return CompleteMultipartResponse(key=key, etag=etag)


async def abort_multipart_upload(bucket: str, key: str, upload_id: str) -> None:
    """Abort a multipart upload, discarding uploaded parts."""
    async with get_s3_client() as s3:
        await s3.abort_multipart_upload(Bucket=bucket, Key=key, UploadId=upload_id)
    logger.info("AbortMultipartUpload: bucket=%s key=%s upload_id=%s", bucket, key, upload_id)


# =========================================================================== #
# Presigned URLs (download / simple upload)
# =========================================================================== #


async def generate_download_url(bucket: str, key: str) -> PresignedUrlResponse:
    """Generate a presigned GET URL for downloading an object."""
    expiry = await _presigned_url_expiry()
    async with get_s3_client() as s3:
        url = await s3.generate_presigned_url(
            "get_object",
            Params={"Bucket": bucket, "Key": key},
            ExpiresIn=expiry,
        )
    return PresignedUrlResponse(url=url, expires_in=expiry)


async def generate_upload_url(
    bucket: str,
    key: str,
    content_type: str = "application/octet-stream",
) -> PresignedUploadUrlResponse:
    """Generate a presigned PUT URL for uploading an object directly."""
    expiry = await _presigned_url_expiry()
    async with get_s3_client() as s3:
        url = await s3.generate_presigned_url(
            "put_object",
            Params={
                "Bucket": bucket,
                "Key": key,
                "ContentType": content_type,
            },
            ExpiresIn=expiry,
        )
    return PresignedUploadUrlResponse(
        url=url,
        key=key,
        expires_in=expiry,
    )


# =========================================================================== #
# 4. S3 Select
# =========================================================================== #


def _build_input_serialization(
    input_format: str,
    compression: str,
    csv_header: str,
) -> dict:
    """Build the InputSerialization dict for S3 Select."""
    ser: dict = {"CompressionType": compression}
    fmt = input_format.upper()
    if fmt == "CSV":
        ser["CSV"] = {"FileHeaderInfo": csv_header}
    elif fmt == "JSON":
        ser["JSON"] = {"Type": "DOCUMENT"}
    elif fmt == "PARQUET":
        ser["Parquet"] = {}
    return ser


def _build_output_serialization(output_format: str) -> dict:
    """Build the OutputSerialization dict for S3 Select."""
    fmt = output_format.upper()
    if fmt == "CSV":
        return {"CSV": {}}
    return {"JSON": {}}


async def s3_select(
    bucket: str,
    key: str,
    expression: str,
    input_format: str = "CSV",
    output_format: str = "JSON",
    compression: str = "NONE",
    csv_header: str = "USE",
) -> S3SelectResponse:
    """Execute an S3 Select SQL query on a CSV/JSON/Parquet object."""
    async with get_s3_client() as s3:
        resp = await s3.select_object_content(
            Bucket=bucket,
            Key=key,
            Expression=expression,
            ExpressionType="SQL",
            InputSerialization=_build_input_serialization(input_format, compression, csv_header),
            OutputSerialization=_build_output_serialization(output_format),
        )

        records: list[str] = []
        stats: dict | None = None

        async for event in resp["Payload"]:
            if "Records" in event:
                payload = event["Records"]["Payload"]
                if isinstance(payload, bytes):
                    records.append(payload.decode("utf-8"))
                else:
                    records.append(str(payload))
            elif "Stats" in event:
                details = event["Stats"].get("Details", {})
                stats = {
                    "bytes_scanned": details.get("BytesScanned"),
                    "bytes_processed": details.get("BytesProcessed"),
                    "bytes_returned": details.get("BytesReturned"),
                }

    logger.info("S3Select: bucket=%s key=%s records=%d", bucket, key, len(records))
    return S3SelectResponse(records=records, stats=stats)


# =========================================================================== #
# 5. Object Tagging
# =========================================================================== #


async def get_object_tagging(bucket: str, key: str) -> TaggingResponse:
    """Get the tag set of an object."""
    async with get_s3_client() as s3:
        resp = await s3.get_object_tagging(Bucket=bucket, Key=key)
    tags = [
        Tag(Key=t["Key"], Value=t["Value"])
        for t in resp.get("TagSet", [])
    ]
    return TaggingResponse(key=key, tags=tags)


async def put_object_tagging(bucket: str, key: str, tags: list[Tag]) -> None:
    """Replace the tag set of an object."""
    tag_set = [{"Key": t.key, "Value": t.value} for t in tags]
    async with get_s3_client() as s3:
        await s3.put_object_tagging(
            Bucket=bucket,
            Key=key,
            Tagging={"TagSet": tag_set},
        )
    logger.info("PutObjectTagging: bucket=%s key=%s tags=%d", bucket, key, len(tags))


async def delete_object_tagging(bucket: str, key: str) -> None:
    """Remove all tags from an object."""
    async with get_s3_client() as s3:
        await s3.delete_object_tagging(Bucket=bucket, Key=key)
    logger.info("DeleteObjectTagging: bucket=%s key=%s", bucket, key)


# =========================================================================== #
# 6. File Preview
# =========================================================================== #

MAX_PREVIEW_ROWS = 1000


async def _download_object_bytes(bucket: str, key: str) -> bytes:
    """Download object bytes from S3."""
    async with get_s3_client() as s3:
        resp = await s3.get_object(Bucket=bucket, Key=key)
        return await resp["Body"].read()


async def preview_parquet(
    bucket: str, key: str, max_rows: int = MAX_PREVIEW_ROWS,
) -> TablePreviewResponse:
    """Preview a Parquet file as tabular data using PyArrow."""
    import io
    import pyarrow.parquet as pq

    data = await _download_object_bytes(bucket, key)
    pf = pq.ParquetFile(io.BytesIO(data))
    total_rows = pf.metadata.num_rows
    columns = pf.schema_arrow.names

    table = pf.read_row_groups(list(range(pf.metadata.num_row_groups)))
    if table.num_rows > max_rows:
        table = table.slice(0, max_rows)

    rows = []
    for batch in table.to_batches():
        cols = [batch.column(i).to_pylist() for i in range(batch.num_columns)]
        for row_idx in range(batch.num_rows):
            rows.append([_serialize_value(cols[ci][row_idx]) for ci in range(len(cols))])

    logger.info("PreviewParquet: bucket=%s key=%s rows=%d/%d", bucket, key, len(rows), total_rows)
    return TablePreviewResponse(
        format="parquet",
        columns=columns,
        rows=rows,
        total_rows=total_rows,
    )


async def preview_xlsx(
    bucket: str,
    key: str,
    sheet: str | None = None,
    max_rows: int = MAX_PREVIEW_ROWS,
) -> TablePreviewResponse:
    """Preview an XLSX/XLS file as tabular data using openpyxl."""
    import io
    import openpyxl

    data = await _download_object_bytes(bucket, key)
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)

    sheet_names = wb.sheetnames
    active_sheet = sheet if sheet and sheet in sheet_names else sheet_names[0]
    ws = wb[active_sheet]

    rows: list[list] = []
    columns: list[str] = []
    for i, row in enumerate(ws.iter_rows(values_only=True)):
        if i == 0:
            columns = [str(c) if c is not None else f"col_{j}" for j, c in enumerate(row)]
            continue
        if i > max_rows:
            break
        rows.append([_serialize_value(c) for c in row])

    total_rows = ws.max_row - 1 if ws.max_row else 0
    wb.close()

    ext = key.rsplit(".", 1)[-1].lower() if "." in key else "xlsx"
    logger.info(
        "PreviewXlsx: bucket=%s key=%s sheet=%s rows=%d/%d",
        bucket, key, active_sheet, len(rows), total_rows,
    )
    return TablePreviewResponse(
        format=ext,
        columns=columns,
        rows=rows,
        total_rows=total_rows,
        sheet_names=sheet_names,
        active_sheet=active_sheet,
    )


async def preview_docx(bucket: str, key: str) -> DocumentPreviewResponse:
    """Preview a DOCX file by converting to HTML using mammoth."""
    import io
    import mammoth

    data = await _download_object_bytes(bucket, key)
    result = mammoth.convert_to_html(io.BytesIO(data))
    if result.messages:
        logger.warning("PreviewDocx warnings: %s", result.messages)

    logger.info("PreviewDocx: bucket=%s key=%s html_len=%d", bucket, key, len(result.value))
    return DocumentPreviewResponse(format="docx", html=result.value)


async def preview_pptx(bucket: str, key: str) -> DocumentPreviewResponse:
    """Preview a PPTX file by extracting slide text and notes."""
    import io
    from pptx import Presentation

    data = await _download_object_bytes(bucket, key)
    prs = Presentation(io.BytesIO(data))

    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(row_texts))

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides.append({
            "slide_number": i,
            "texts": texts,
            "notes": notes,
        })

    # Build a simple HTML representation
    html_parts = []
    for s in slides:
        html_parts.append(f'<div class="slide"><h3>Slide {s["slide_number"]}</h3>')
        for t in s["texts"]:
            html_parts.append(f"<p>{t}</p>")
        if s["notes"]:
            html_parts.append(f'<blockquote class="notes">{s["notes"]}</blockquote>')
        html_parts.append("</div><hr/>")

    logger.info("PreviewPptx: bucket=%s key=%s slides=%d", bucket, key, len(slides))
    return DocumentPreviewResponse(
        format="pptx",
        html="\n".join(html_parts),
        slides=slides,
    )


# =========================================================================== #
# 7. File Browser Configuration
# =========================================================================== #


async def get_filebrowser_config(session: AsyncSession) -> FilebrowserConfigResponse:
    """Load File Browser configuration from the database.

    Queries three tables and assembles a single response:
      - argus_configuration_filebrowser → browser-level key-value settings
      - argus_configuration_filebrowser_preview → per-category preview limits
      - argus_configuration_filebrowser_extension → extension-to-category mapping
    """
    # 1) Browser global settings
    result = await session.execute(select(ArgusConfigurationFilebrowser))
    browser_rows = result.scalars().all()
    browser: dict[str, int] = {}
    for row in browser_rows:
        try:
            browser[row.config_key] = int(row.config_value)
        except ValueError:
            browser[row.config_key] = row.config_value

    # 2) Preview categories
    result = await session.execute(select(ArgusConfigurationFilebrowserPreview))
    preview_rows = result.scalars().all()

    # 3) Extensions grouped by preview_id
    result = await session.execute(select(ArgusConfigurationFilebrowserExtension))
    ext_rows = result.scalars().all()
    ext_map: dict[int, list[str]] = defaultdict(list)
    for ext in ext_rows:
        ext_map[ext.preview_id].append(ext.extension)

    # Assemble response
    preview: list[PreviewCategoryResponse] = []
    for cat in preview_rows:
        extensions = sorted(ext_map.get(cat.id, []))
        preview.append(
            PreviewCategoryResponse(
                category=cat.category,
                label=cat.label,
                extensions=extensions,
                max_file_size=cat.max_file_size,
                max_preview_rows=cat.max_preview_rows,
            )
        )

    logger.info("FilebrowserConfig: browser_keys=%d categories=%d", len(browser), len(preview))
    return FilebrowserConfigResponse(browser=browser, preview=preview)


async def update_browser_settings(session: AsyncSession, settings_map: dict[str, int]) -> None:
    """Update browser-level key-value settings."""
    for key, value in settings_map.items():
        result = await session.execute(
            select(ArgusConfigurationFilebrowser).where(
                ArgusConfigurationFilebrowser.config_key == key
            )
        )
        row = result.scalar_one_or_none()
        if row:
            row.config_value = str(value)
        else:
            session.add(ArgusConfigurationFilebrowser(
                config_key=key, config_value=str(value),
            ))
    await session.commit()
    logger.info("UpdateBrowserSettings: keys=%s", list(settings_map.keys()))


async def update_preview_category(
    session: AsyncSession,
    category: str,
    max_file_size: int,
    max_preview_rows: int | None,
) -> None:
    """Update a single preview category's limits."""
    result = await session.execute(
        select(ArgusConfigurationFilebrowserPreview).where(
            ArgusConfigurationFilebrowserPreview.category == category
        )
    )
    row = result.scalar_one_or_none()
    if not row:
        raise ValueError(f"Preview category not found: {category}")
    row.max_file_size = max_file_size
    row.max_preview_rows = max_preview_rows
    await session.commit()
    logger.info(
        "UpdatePreviewCategory: category=%s max_file_size=%d max_preview_rows=%s",
        category, max_file_size, max_preview_rows,
    )


def _serialize_value(val):
    """Convert a cell value to a JSON-safe type."""
    if val is None:
        return None
    if isinstance(val, (int, float, str, bool)):
        return val
    if isinstance(val, datetime):
        return val.isoformat()
    if isinstance(val, bytes):
        return val.hex()
    return str(val)
