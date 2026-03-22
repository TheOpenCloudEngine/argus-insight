"""Local filesystem browser service.

Provides file browsing operations on the local Linux filesystem using
os/pathlib. All paths are resolved relative to the configured data directory
(settings.data_dir). When the API receives path="/", it maps to data_dir.
"""

import grp
import logging
import os
import pwd
import shutil
import stat
from datetime import UTC, datetime
from pathlib import Path

from app.core.config import settings
from app.filesystemmgr.schemas import (
    CreateFolderResponse,
    DeleteResponse,
    DocumentPreviewResponse,
    FileInfo,
    FileStatResponse,
    FolderInfo,
    ListDirectoryResponse,
    RenameResponse,
    TablePreviewResponse,
)

logger = logging.getLogger(__name__)


# =========================================================================== #
# Helpers
# =========================================================================== #


def _get_root_dir(sub_path: str | None = None) -> Path:
    """Return the root directory for the file browser.

    If sub_path is provided, returns data_dir / sub_path.
    e.g. sub_path="model-artifacts" → data_dir/model-artifacts
    """
    root = settings.data_dir
    if sub_path:
        root = root / sub_path
    return root


def _format_dt(timestamp: float) -> str:
    """Convert a Unix timestamp to ISO 8601 string."""
    return datetime.fromtimestamp(timestamp, tz=UTC).isoformat()


def _permission_string(mode: int) -> str:
    """Convert a stat mode to a rwxrwxrwx permission string."""
    perms = ""
    for who in ("USR", "GRP", "OTH"):
        for what in ("R", "W", "X"):
            flag = getattr(stat, f"S_I{what}{who}")
            perms += what.lower() if mode & flag else "-"
    return perms


def _permission_octal(mode: int) -> str:
    """Convert a stat mode to octal string (e.g. '0755')."""
    return format(stat.S_IMODE(mode), "04o")


def _get_owner(st: os.stat_result) -> str:
    """Get username for a file's uid."""
    try:
        return pwd.getpwuid(st.st_uid).pw_name
    except KeyError:
        return str(st.st_uid)


def _get_group(st: os.stat_result) -> str:
    """Get group name for a file's gid."""
    try:
        return grp.getgrgid(st.st_gid).gr_name
    except KeyError:
        return str(st.st_gid)


def _resolve_path(path: str, root_sub: str | None = None) -> Path:
    """Resolve a user-supplied path relative to root directory.

    The user sees "/" as the root, which maps to data_dir (or data_dir/root_sub).
    For example, if data_dir is /var/lib/argus-catalog-server and root_sub="model-artifacts":
      - path="/"         -> /var/lib/argus-catalog-server/model-artifacts
      - path="/iris"     -> /var/lib/argus-catalog-server/model-artifacts/iris
    """
    root = _get_root_dir(root_sub).resolve()
    root.mkdir(parents=True, exist_ok=True)

    relative = path.lstrip("/")
    if relative:
        resolved = (root / relative).resolve()
    else:
        resolved = root

    if not (resolved == root or str(resolved).startswith(str(root) + "/")):
        raise ValueError("Access denied: path escapes data directory")

    return resolved


def _to_user_path(resolved: Path, root_sub: str | None = None) -> str:
    """Convert an internal resolved path back to a user-visible path."""
    root = _get_root_dir(root_sub).resolve()
    try:
        rel = resolved.relative_to(root)
        return "/" + str(rel) if str(rel) != "." else "/"
    except ValueError:
        return str(resolved)


# =========================================================================== #
# 1. Directory listing
# =========================================================================== #


async def list_directory(path: str, root_sub: str | None = None) -> ListDirectoryResponse:
    """List files and directories under a given path."""
    resolved = _resolve_path(path, root_sub)

    if not resolved.exists():
        raise FileNotFoundError(f"Directory not found: {resolved}")
    if not resolved.is_dir():
        raise NotADirectoryError(f"Not a directory: {resolved}")

    folders: list[FolderInfo] = []
    files: list[FileInfo] = []

    try:
        entries = sorted(resolved.iterdir(), key=lambda e: e.name)
    except PermissionError:
        raise PermissionError(f"Permission denied: {resolved}")

    for entry in entries:
        try:
            st = entry.stat(follow_symlinks=False)
        except (PermissionError, OSError):
            continue

        owner = _get_owner(st)
        group = _get_group(st)
        perms = _permission_string(st.st_mode)

        if stat.S_ISDIR(st.st_mode):
            folders.append(FolderInfo(
                key=_to_user_path(entry, root_sub) + "/",
                name=entry.name,
                owner=owner,
                group=group,
                permissions=perms,
            ))
        else:
            files.append(FileInfo(
                key=_to_user_path(entry, root_sub),
                name=entry.name,
                size=st.st_size,
                last_modified=_format_dt(st.st_mtime),
                owner=owner,
                group=group,
                permissions=perms,
            ))

    return ListDirectoryResponse(
        folders=folders,
        files=files,
        current_path=_to_user_path(resolved, root_sub),
    )


# =========================================================================== #
# 2. Create directory
# =========================================================================== #


async def create_folder(path: str, root_sub: str | None = None) -> CreateFolderResponse:
    """Create a new directory (including parents)."""
    resolved = _resolve_path(path, root_sub)
    resolved.mkdir(parents=True, exist_ok=True)
    logger.info("CreateFolder: %s", resolved)
    return CreateFolderResponse(path=_to_user_path(resolved, root_sub))


# =========================================================================== #
# 3. Delete files/directories
# =========================================================================== #


async def delete_paths(paths: list[str], root_sub: str | None = None) -> DeleteResponse:
    """Delete multiple files or directories."""
    deleted: list[str] = []
    errors: list[dict] = []

    for p in paths:
        try:
            resolved = _resolve_path(p, root_sub)
            if not resolved.exists():
                errors.append({"path": p, "error": "Not found"})
                continue

            if resolved.is_dir():
                shutil.rmtree(resolved)
            else:
                resolved.unlink()

            deleted.append(p)
            logger.info("Deleted: %s", resolved)
        except Exception as e:
            errors.append({"path": p, "error": str(e)})
            logger.error("Delete error: %s - %s", p, e)

    return DeleteResponse(deleted=deleted, errors=errors)


# =========================================================================== #
# 4. Rename / Move
# =========================================================================== #


async def rename(source_path: str, destination_path: str, root_sub: str | None = None) -> RenameResponse:
    """Rename or move a file/directory."""
    src = _resolve_path(source_path, root_sub)
    dst = _resolve_path(destination_path, root_sub)

    if not src.exists():
        raise FileNotFoundError(f"Source not found: {src}")
    if dst.exists():
        raise FileExistsError(f"Destination already exists: {dst}")

    src.rename(dst)
    logger.info("Rename: %s -> %s", src, dst)
    return RenameResponse(source=_to_user_path(src, root_sub), destination=_to_user_path(dst, root_sub))


# =========================================================================== #
# 5. File metadata (stat)
# =========================================================================== #


async def file_stat(path: str, root_sub: str | None = None) -> FileStatResponse:
    """Get detailed metadata for a file or directory."""
    resolved = _resolve_path(path, root_sub)

    if not resolved.exists():
        raise FileNotFoundError(f"Not found: {resolved}")

    st = resolved.stat(follow_symlinks=False)
    symlink_target = str(resolved.readlink()) if resolved.is_symlink() else None

    return FileStatResponse(
        path=_to_user_path(resolved, root_sub),
        name=resolved.name or "/",
        is_directory=resolved.is_dir(),
        size=st.st_size,
        last_modified=_format_dt(st.st_mtime),
        last_accessed=_format_dt(st.st_atime),
        created=_format_dt(st.st_ctime),
        owner=_get_owner(st),
        group=_get_group(st),
        permissions=_permission_string(st.st_mode),
        permissions_octal=_permission_octal(st.st_mode),
        inode=st.st_ino,
        hard_links=st.st_nlink,
        symlink_target=symlink_target,
    )


# =========================================================================== #
# 6. File download (read bytes)
# =========================================================================== #


async def read_file(path: str, root_sub: str | None = None) -> tuple[bytes, str]:
    """Read a file and return (bytes, filename).

    Returns the raw bytes so the router can build a streaming response.
    """
    resolved = _resolve_path(path, root_sub)

    if not resolved.exists():
        raise FileNotFoundError(f"File not found: {resolved}")
    if resolved.is_dir():
        raise IsADirectoryError(f"Cannot download a directory: {resolved}")

    data = resolved.read_bytes()
    return data, resolved.name


# =========================================================================== #
# 7. File upload
# =========================================================================== #


async def save_uploaded_file(
    destination_dir: str,
    filename: str,
    content: bytes,
    root_sub: str | None = None,
) -> str:
    """Save an uploaded file to the specified directory."""
    dir_path = _resolve_path(destination_dir, root_sub)
    if not dir_path.is_dir():
        raise NotADirectoryError(f"Not a directory: {dir_path}")

    file_path = dir_path / filename
    file_path.write_bytes(content)
    logger.info("Upload: %s (%d bytes)", file_path, len(content))
    return _to_user_path(file_path, root_sub)


# =========================================================================== #
# 8. File Preview
# =========================================================================== #

MAX_PREVIEW_ROWS = 1000


def _serialize_value(val):
    """Convert a cell value to a JSON-safe type."""
    if val is None:
        return None
    if isinstance(val, (int, float, str, bool)):
        return val
    if isinstance(val, datetime):
        return val.isoformat()
    if isinstance(val, bytes):
        return val.hex()
    return str(val)


async def preview_parquet(path: str, max_rows: int = MAX_PREVIEW_ROWS, root_sub: str | None = None) -> TablePreviewResponse:
    """Preview a Parquet file as tabular data using PyArrow."""
    import io
    import pyarrow.parquet as pq

    resolved = _resolve_path(path, root_sub)
    data = resolved.read_bytes()
    pf = pq.ParquetFile(io.BytesIO(data))
    total_rows = pf.metadata.num_rows
    columns = pf.schema_arrow.names

    table = pf.read_row_groups(list(range(pf.metadata.num_row_groups)))
    if table.num_rows > max_rows:
        table = table.slice(0, max_rows)

    rows = []
    for batch in table.to_batches():
        cols = [batch.column(i).to_pylist() for i in range(batch.num_columns)]
        for row_idx in range(batch.num_rows):
            rows.append([_serialize_value(cols[ci][row_idx]) for ci in range(len(cols))])

    logger.info("PreviewParquet: %s rows=%d/%d", resolved, len(rows), total_rows)
    return TablePreviewResponse(
        format="parquet",
        columns=columns,
        rows=rows,
        total_rows=total_rows,
    )


async def preview_xlsx(
    path: str,
    sheet: str | None = None,
    max_rows: int = MAX_PREVIEW_ROWS,
    root_sub: str | None = None,
) -> TablePreviewResponse:
    """Preview an XLSX/XLS file as tabular data using openpyxl."""
    import io
    import openpyxl

    resolved = _resolve_path(path, root_sub)
    data = resolved.read_bytes()
    wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)

    sheet_names = wb.sheetnames
    active_sheet = sheet if sheet and sheet in sheet_names else sheet_names[0]
    ws = wb[active_sheet]

    rows: list[list] = []
    columns: list[str] = []
    for i, row in enumerate(ws.iter_rows(values_only=True)):
        if i == 0:
            columns = [str(c) if c is not None else f"col_{j}" for j, c in enumerate(row)]
            continue
        if i > max_rows:
            break
        rows.append([_serialize_value(c) for c in row])

    total_rows = ws.max_row - 1 if ws.max_row else 0
    wb.close()

    ext = resolved.suffix.lstrip(".").lower() or "xlsx"
    logger.info("PreviewXlsx: %s sheet=%s rows=%d/%d", resolved, active_sheet, len(rows), total_rows)
    return TablePreviewResponse(
        format=ext,
        columns=columns,
        rows=rows,
        total_rows=total_rows,
        sheet_names=sheet_names,
        active_sheet=active_sheet,
    )


async def preview_docx(path: str, root_sub: str | None = None) -> DocumentPreviewResponse:
    """Preview a DOCX file by converting to HTML using mammoth."""
    import io
    import mammoth

    resolved = _resolve_path(path, root_sub)
    data = resolved.read_bytes()
    result = mammoth.convert_to_html(io.BytesIO(data))
    if result.messages:
        logger.warning("PreviewDocx warnings: %s", result.messages)

    logger.info("PreviewDocx: %s html_len=%d", resolved, len(result.value))
    return DocumentPreviewResponse(format="docx", html=result.value)


async def preview_pptx(path: str, root_sub: str | None = None) -> DocumentPreviewResponse:
    """Preview a PPTX file by extracting slide text and notes."""
    import io
    from pptx import Presentation

    resolved = _resolve_path(path, root_sub)
    data = resolved.read_bytes()
    prs = Presentation(io.BytesIO(data))

    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells]
                    texts.append(" | ".join(row_texts))

        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        slides.append({
            "slide_number": i,
            "texts": texts,
            "notes": notes,
        })

    html_parts = []
    for s in slides:
        html_parts.append(f'<div class="slide"><h3>Slide {s["slide_number"]}</h3>')
        for t in s["texts"]:
            html_parts.append(f"<p>{t}</p>")
        if s["notes"]:
            html_parts.append(f'<blockquote class="notes">{s["notes"]}</blockquote>')
        html_parts.append("</div><hr/>")

    logger.info("PreviewPptx: %s slides=%d", resolved, len(slides))
    return DocumentPreviewResponse(
        format="pptx",
        html="\n".join(html_parts),
        slides=slides,
    )
